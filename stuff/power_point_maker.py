from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os
import datetime

# Updated path to the folder containing images
image_folder_path = r'C:\Users\sapir\Pictures\Screenshots\images2'

# Create a PowerPoint presentation
presentation = Presentation()

# Iterate through the images in the folder
for filename in os.listdir(image_folder_path):
    if filename.endswith(('.png', '.jpg', '.jpeg')):
        # Create a new slide for each image
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])

        # Load the image and convert to RGB mode, then crop to keep only the left corner
        image_path = os.path.join(image_folder_path, filename)
        img = Image.open(image_path).convert('RGB')
        width, height = img.size
        x= 50
        y= 50
        left_corner = img.crop((0+x+140, 0+y, width // 2-20+0, height // 2+2*y+155))

        # Save the cropped image temporarily
        temp_image_path = os.path.join(image_folder_path, 'temp_cropped_image.jpg')
        left_corner.save(temp_image_path)

        # Move the image a little to the right and down
        left = Inches(1)  # Adjust the left position as needed
        top = Inches(1)   # Adjust the top position as needed

        # Add the cropped image to the slide and resize it to fit the slide
        pic = slide.shapes.add_picture(temp_image_path, left, top, width=Inches(8), height=Inches(5))

# Remove the temporary cropped image
os.remove(temp_image_path)

# Generate the output file name with the current date and time
current_datetime = datetime.datetime.now().strftime('%Y%m%d_%H%M%S')
output_file_path = r'C:\Users\sapir\Pictures\Screenshots\images2\ImagePresentation2_' + current_datetime + '.pptx'

# Save the presentation
presentation.save(output_file_path)


"""
main
from pptx import Presentation
from pptx.util import Inches
import os

# Path to the folder containing images
image_folder_path = r'C:\Users\sapir\Pictures\Screenshots\images'

# Create a PowerPoint presentation
presentation = Presentation()

# Iterate through the images in the folder
for filename in os.listdir(image_folder_path):
    if filename.endswith(('.png', '.jpg', '.jpeg')):
        # Create a new slide for each image
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])

        # Add the image to the slide
        image_path = os.path.join(image_folder_path, filename)
        left = Inches(1)
        top = Inches(1)
        pic = slide.shapes.add_picture(image_path, left, top)

# Save the presentation
presentation.save('ImagePresentation.pptx')


"""


"""
fit


from pptx import Presentation
from pptx.util import Inches
import os

# Path to the folder containing images
image_folder_path = r'C:\Users\sapir\Pictures\Screenshots\images'

# Create a PowerPoint presentation
presentation = Presentation()

# Iterate through the images in the folder
for filename in os.listdir(image_folder_path):
    if filename.endswith(('.png', '.jpg', '.jpeg')):
        # Create a new slide for each image
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])

        # Add the image to the slide and resize it to fit the slide
        image_path = os.path.join(image_folder_path, filename)
        left = top = Inches(1)
        pic = slide.shapes.add_picture(image_path, left, top, width=Inches(8), height=Inches(5))

# Save the presentation
presentation.save('ImagePresentation.pptx')


"""